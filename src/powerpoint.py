import collections.abc # インポートしないとエラーが発生する
from pptx.util import Inches  # インチ
import pptx

from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.util import Cm,Pt               # 単位指定をするクラス(センチメートル, ポイント単位)
#from pptx.enum.text import PP_ALIGN  # 中央揃えにする用
#from pptx import Presentation # プレゼンテーションを作成
#import os

#cd = os.getcwd()
prs = pptx.Presentation()
prs.slide_width = Inches(11.69) #A4サイズ
prs.slide_height = Inches(8.27)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白のスライドを追加



start_height = 0.5

for i in range(2):
    start_left = 1
    for j in range(4):
        pic_left = Cm(start_left)
        pic_top = Cm(start_height)

        pic_height = Cm(9.5)

        image_path = f'./input/IMG_{j:03d}.png'
        slide.shapes.add_picture(image_path, pic_left, pic_top, height=pic_height)

        rect0 = slide.shapes.add_shape(		# shapeオブジェクト➀を追加
        MSO_SHAPE.RIGHT_ARROW,   	                    # 図形の種類を[丸角四角形]に指定
        Cm(3), Cm(2),               # 挿入位置の指定　左上の座標の指定
        Cm(5), Cm(3))               # 挿入図形の幅と高さの指定

        start_left += 7
    
    start_height += 10




prs.save("./test.pptx")