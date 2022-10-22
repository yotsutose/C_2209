# tkinter
import tkinter
import tkinter.filedialog
import tkinter.ttk as ttk

# 画像処理
from PIL import Image, ImageTk
import cv2
import numpy as np

# pptx
import collections.abc
from operator import imod # インポートしないとエラーが発生する
import pptx
from pptx.enum.shapes import MSO_SHAPE # 図形の定義がされているクラス
from pptx.dml.color import RGBColor
from pptx.util import Cm,Pt,Inches # 単位指定をするクラス(センチメートル, ポイント単位)
from pptx.enum.text import PP_ALIGN # 段落の水平位置のEnume

# 仮想ファイルシステム
import tempfile
import os

# プログラマ便利ツール
from tqdm import tqdm # プログレスバーの導入

class Model():

    def __init__(self):

        # 動画オブジェクト参照用
        self.cap= None

        # 読み込んだフレーム
        self.frames = []

        # mode 0:プレビュー、1:編集
        self.mode = tkinter.IntVar()

        # 読み込んだフレームの選択状態 1:選択、0:未選択
        # 最初から2は消しておくかもしれない
        # 0と1しかない予定
        self.frame_state = []

        # PIL画像オブジェクト参照用
        self.image = None

        # Tkinter画像オブジェクト参照用
        self.image_tk = None

        # 現在表示中のフレーム
        self.now = tkinter.IntVar()
        self.nows = []
        for x in range(5):
            self.nows.append(tkinter.IntVar())
            self.nows[x].set(-1)

        self.now.trace_add("write", self.set_nows)

        # debug用に最初から読み込んでおく
        self.create_video("./input/input1.MP4")

    def create_video(self, path):
        '動画オブジェクトの生成を行う'

        # pathの動画から動画オブジェクト生成
        self.cap = cv2.VideoCapture(path)
        width = int(self.cap.get(cv2.CAP_PROP_FRAME_WIDTH))
        height = int(self.cap.get(cv2.CAP_PROP_FRAME_HEIGHT))
        frame_count = int(self.cap.get(cv2.CAP_PROP_FRAME_COUNT))
        print(f'{width=} {height=} {frame_count=}')

        self.frames = [] # 1/4スケールの画像保存
        self.perfect_frames = [] # 劣化なしの状態
        self.frame_state = [] # 0か1で状態を書く

        self.first_img = True
        self.pre_img_is_similar = False
        self.diff_max = 5
        
        print(f'ビデオの読み込み中')
        for i in tqdm(range(frame_count)):
            ret, image2 = self.cap.read()
            if ret == False:
                break
            if i % 5 != 0: # 5枚ずつ切り取る
                continue
            if i == 0:
                image1 = image2
                continue
            
            height = image1.shape[0]
            width = image1.shape[1]
            img_size = (int(width), int(height))
            ret = self.func5(image1, image2, img_size)

            if ret == 1:
                img = image1
            elif ret == 2:
                img = image2
            else:
                image1 = image2
                continue

            # 完璧なイメージを保存 意外と時間はかからない
            self.perfect_frames.append(img)

            # プレビューの表示用に保存
            rgb_frame = cv2.cvtColor(img, cv2.COLOR_BGR2RGB)
            pil_image = Image.fromarray(rgb_frame)
            pil_image = pil_image.resize((width//4, height//4), resample=3)
            self.frames.append(ImageTk.PhotoImage(pil_image))

            # 状態管理用に追加
            self.frame_state.append(tkinter.IntVar())
            self.frame_state[-1].set(ret%2)

            # image1に保存する
            image1 = image2
        # セットする
        self.now.set(0)

    def func5(self, image1, image2, img_size):
        # 比較するために、同じサイズにリサイズしておく
        image1 = cv2.resize(image1, img_size)
        image2 = cv2.resize(image2, img_size)

        diff = image1.astype(int) - image2.astype(int)
        diff_abs = np.abs(diff)
        degree_of_similarity = np.count_nonzero(diff_abs < self.diff_max) / image2.size

        # 最後の画像を保存
        if(degree_of_similarity > 0.85 and not(self.pre_img_is_similar)): # 同じやつになり始めた
            self.pre_img_is_similar = True
            return 2
        elif(degree_of_similarity <= 0.85 and self.pre_img_is_similar): # 変わり始めた
            self.pre_img_is_similar = False
            return 1
        return 0

    def set_nows(self, a, b, c):
        for x in range(5):
            self.nows[x].set(-1)
        if self.mode.get() == 0: # プレビューモード
            now = self.now.get()
            idx = 0
            for i in range(now, len(self.frames)-5, 1):
                if self.frame_state[i].get() == 1:
                    self.nows[idx].set(i)
                    idx += 1
                if idx >= 5:
                    break
        else: # 編集モード
            for i in range(5):
                state_ = self.now.get() + i
                if len(self.frames) <= state_:
                    continue
                self.nows[i].set(state_)

    def next_frame(self):
        next = min(self.now.get()+1, len(self.frames)-5)
        self.now.set(next)

    def previous_frame(self):
        next = max(self.now.get()-1, 0)
        self.now.set(next)


class View():

    def __init__(self, app, model):

        self.master = app
        self.model = model

        # callbackを用意
        # for x in range(5):
        #     self.model.nows[x].trace_add("write", lambda name, index, mode: self.make_draw_image(self, x))

        self.model.nows[0].trace_add("write", lambda name, index, mode: self.draw_image(self, 0))
        self.model.nows[1].trace_add("write", lambda name, index, mode: self.draw_image(self, 1))
        self.model.nows[2].trace_add("write", lambda name, index, mode: self.draw_image(self, 2))
        self.model.nows[3].trace_add("write", lambda name, index, mode: self.draw_image(self, 3))
        self.model.nows[4].trace_add("write", lambda name, index, mode: self.draw_image(self, 4))

        # アプリ内のウィジェットを作成
        self.create_widgets()

    def create_widgets(self):
        'アプリ内にウィジェットを作成・配置する'

        self.master.grid_rowconfigure(0, weight=1)
        self.master.grid_columnconfigure(0, weight=1)

        # メインフレーム(編集とプレビュー画面)
        self.main_frame = tkinter.Frame(
            self.master,
            bg="#FCFFEE"
        )
        self.main_frame.grid(row=0, column=0, sticky="nsew")
        self.create_main_frame()

        # ホームフレーム(ホーム画面)
        self.home_frame = tkinter.Frame(
            self.master,
            bg="#FCFFEE"
        )
        self.home_frame.grid(row=0, column=0, sticky="nsew")
        self.home_frame.grid_rowconfigure(0, weight=4)
        self.home_frame.grid_columnconfigure(0, weight=1)
        self.create_home_frame()

        # エンドフレーム(終了画面)
        self.end_frame = tkinter.Frame(
            self.master,
            bg="#FCFFEE"
        )
        self.end_frame.grid(row=0, column=0, sticky="nsew")
        self.end_frame.grid_rowconfigure(0, weight=1)
        self.end_frame.grid_columnconfigure(0, weight=1)
        self.create_end_frame()

        # tkraise()
        # self.home_frame.tkraise()
        self.main_frame.tkraise()
        # self.end_frame.tkraise()


    def create_main_frame(self):
        # 画像の読み込み
        edit_top_image_ = Image.open('./app/image/編集モード.png')
        self.edit_top_image = ImageTk.PhotoImage(edit_top_image_)

        preview_top_image_ = Image.open('./app/image/プレビューモード.png')
        self.preview_top_image = ImageTk.PhotoImage(preview_top_image_)

        add_button_image_ = Image.open('./app/image/未選択_四角.png')
        self.add_button_image = ImageTk.PhotoImage(add_button_image_)

        delete_button_image_ = Image.open('./app/image/選択_四角.png')
        self.delete_button_image = ImageTk.PhotoImage(delete_button_image_)

        next_frame_button_image_ = Image.open('./app/image/次へ_Y.png')
        self.next_frame_button_image = ImageTk.PhotoImage(next_frame_button_image_)
        
        prev_frame_button_image_ = Image.open('./app/image/前へ_M.png')
        self.prev_frame_button_image = ImageTk.PhotoImage(prev_frame_button_image_)
        
        mode_button_image_ = Image.open('./app/image/編集.png')
        self.mode_button_image = ImageTk.PhotoImage(mode_button_image_)
        
        mode_button_image2_ = Image.open('./app/image/完了ボタン.png')
        self.mode_button_image2 = ImageTk.PhotoImage(mode_button_image2_)

        to_home_button_image_ = Image.open('./app/image/ホームへ.png')
        self.to_home_button_image = ImageTk.PhotoImage(to_home_button_image_)

        making_pptx_button_image_ = Image.open('./app/image/実行ボタン.png')
        self.making_pptx_button_image = ImageTk.PhotoImage(making_pptx_button_image_)

        # ホームへ/編集モード/プレビューモードへのフレーム
        self.head_frame = tkinter.Frame(
            self.main_frame,
            highlightbackground='#FCFFEE',
            bg="#FCFFEE"
        )
        self.head_frame.grid(column=1, row=1)

        # ホームへのボタン
        self.to_home_button = tkinter.Button(
            self.main_frame,
            image = self.to_home_button_image,
        )
        self.to_home_button.grid(in_=self.head_frame, column=0, row=0)

        # 編集モード/プレビューモードのキャンバス
        self.head_canvas = tkinter.Canvas(
            self.head_frame,
            height=150,
            width=500,
            highlightbackground='#FCFFEE',
            bg="#FCFFEE"
        )
        # 最初はプレビューモードの表示をしておく
        self.head_canvas.create_image(
            30, 30,
            image=self.preview_top_image,
            anchor=tkinter.NW,
            tag="image",
        )
        self.head_canvas.grid(column=1, row=0, padx=150)

        self.to_home_button2 = tkinter.Button(
            self.main_frame,
            image = self.to_home_button_image,
        )
        self.to_home_button2.grid(in_=self.head_frame,column=2, row=0)
        self.to_home_button2.lower()

        # キャンバス5枚のフレーム
        self.canvas_frame = tkinter.Frame(
            self.main_frame,
            bg="#FCFFEE"
        )
        self.canvas_frame.grid(column=1, row=2, padx=100)

        # オペレーションフレーム
        self.operation_frame = tkinter.Frame(
            self.main_frame,
            bg="#FCFFEE"
        )
        self.operation_frame.grid(column=1, row=3)

        # フッターのフレーム
        self.foot_frame = tkinter.Frame(
            self.main_frame,
            height=20,
            width=700,
            bg="#FCFFEE"
        )
        self.foot_frame.grid(column=1, row=4)

        # パネルごとのフレーム in canvas_panels
        self.canvas_paneles = [tkinter.Frame(
            self.canvas_frame,
            bg="#FCFFEE",) for x in range(5)]
        [self.canvas_paneles[x].grid(column=x, row=1, padx=10) for x in range(5)]

        # フレーム番号表示 in canvas_panels
        self.frame_index = [tkinter.Label(
            self.canvas_paneles[x],
            bg="#FCFFEE",
            font=("MSゴシック", "30", "bold"),
            textvariable=self.model.nows[x]) for x in range(5)]
        [self.frame_index[x].grid(row=0, column=0, sticky="nsew") for x in range(5)]

        # フレーム表示 in canvas_panels
        self.frame = [tkinter.Canvas(
            self.canvas_paneles[x],
            width=200, # ここは要調整
            height=450,
            highlightbackground='#FCFFEE',
            bg='#FCFFEE'
            ) for x in range(5)]
        [self.frame[x].grid(row=1, column=0, sticky="nsew") for x in range(5)]

        # ボタンのためのフレーム
        self.button_frame = [tkinter.Frame(
            self.canvas_paneles[x],
            highlightbackground='#FCFFEE',
            bg='#FCFFEE'
            ) for x in range(5)]
        [self.button_frame[x].grid_rowconfigure(0, weight=1) for x in range(5)]
        [self.button_frame[x].grid_columnconfigure(0, weight=1) for x in range(5)]
        [self.button_frame[x].grid(row=2, column=0, sticky="nsew") for x in range(5)]

        # 追加ボタン表示
        self.state_button = [tkinter.Button(
            self.button_frame[x],
            image = self.add_button_image,
            ) for x in range(5)]
        [self.state_button[x].grid(row=0, column=0) for x in range(5)]

        # 削除ボタン表示
        self.state_button2 = [tkinter.Button(
            self.button_frame[x],
            image = self.delete_button_image,
            ) for x in range(5)]
        [self.state_button2[x].grid(row=0, column=0) for x in range(5)]

        [self.state_button[x].tkraise() for x in range(5)]

        # シークバーの表示
        style = ttk.Style()
        style.configure(
            "Horizontal.TScale",
        )
        self.scale_bar = ttk.Scale(
            self.operation_frame,
            style="TScale",
            variable=self.model.now,
            orient="horizontal",
            length=500,
            from_=0,
            to=len(self.model.frames)-5,
            # command=lambda e: self.draw_image()
        )
        self.scale_bar.pack(pady=25)


        # 編集ボタンは一人きり
        self.mode_button = tkinter.Button(
            self.main_frame,
            image=self.mode_button_image,
            bg='#FCFFEE',
            highlightbackground='#FCFFEE'
        )
        self.mode_button.pack(in_=self.operation_frame, fill = 'x', padx=20, side = 'left')

        # prev_frameへのボタン
        self.prev_frame_button = tkinter.Button(
            self.operation_frame,
            image = self.prev_frame_button_image,
        )
        self.prev_frame_button.pack(fill = 'x', padx=10, side = 'left')

        # modeチェンジのためのフレーム
        self.mode_change_frame = tkinter.Frame(
            self.operation_frame,
            highlightbackground='#FCFFEE',
            bg='#FCFFEE')
        self.mode_change_frame.grid_rowconfigure(0, weight=1)
        self.mode_change_frame.grid_columnconfigure(0, weight=1)
        self.mode_change_frame.pack(fill = 'x', padx=20, side = 'right')

        # next_frameへのボタン
        self.next_frame_button = tkinter.Button(
            self.operation_frame,
            image = self.next_frame_button_image,
        )
        self.next_frame_button.pack(fill = 'x', padx=10, side = 'right')

        # 完了ボタン
        self.mode_button2 = tkinter.Button(
            self.operation_frame,
            image=self.mode_button_image2,
            bg='#FCFFEE',
            highlightbackground='#FCFFEE'
        )
        self.mode_button2.grid(in_=self.mode_change_frame, row=0, column=0, sticky="nsew")        
        
        # 実行ボタン - pptxを作る
        self.making_pptx_button = tkinter.Button(
            self.operation_frame,
            image = self.making_pptx_button_image,
        )
        self.making_pptx_button.grid(in_=self.mode_change_frame, row=0, column=0, sticky="nsew")
        self.making_pptx_button.tkraise()
        return
    
    def create_home_frame(self):
        # 画像の読み込み
        title_logo_image_ = Image.open('./app/image/bigtitle.png')
        self.title_logo_image = ImageTk.PhotoImage(title_logo_image_)

        input_button_image_ = Image.open('./app/image/動画選択.png')
        self.input_button_image = ImageTk.PhotoImage(input_button_image_)

        using_button_image_ = Image.open('./app/image/使い方.png')
        self.using_button_image = ImageTk.PhotoImage(using_button_image_)

        done_button_image_ = Image.open('./app/image/実行ボタン.png')
        self.done_button_image = ImageTk.PhotoImage(done_button_image_)

        # 下の余白のためのフレーム
        self.home_head_frame = tkinter.Frame(
            self.home_frame,
            height=50,
            width=500,
            bg="#FCFFEE"
        )
        self.home_head_frame.grid(column=0, row=0)

        # タイトルロゴを配置するフレーム
        self.home_title_canvas = tkinter.Canvas(
            self.home_frame,
            height=280,
            width=850,
            highlightbackground='#FCFFEE',
            bg="#FCFFEE"
        )
        self.home_title_canvas.create_image(
            30, 80,
            image=self.title_logo_image,
            anchor=tkinter.NW,
            tag="image",
        )
        self.home_title_canvas.grid(column=0, row=1, pady=10)

        # ファイル読み込みを配置するフレーム
        self.home_input_frame = tkinter.Frame(
            self.home_frame,
            height=200,
            width=500,
            highlightbackground='#FCFFEE',
            bg="#FCFFEE"
        )
        self.home_input_frame.grid(column=0, row=2, pady=10)

        # ユーザ操作用フレームの作成と配置
        self.home_operation_frame = tkinter.Frame(
            self.home_frame,
            height=150,
            width=500,
            bg="#FCFFEE"
        )
        self.home_operation_frame.grid(column=0, row=3, pady=10)

        # 下の余白のためのフレーム
        self.home_foot_frame = tkinter.Frame(
            self.home_frame,
            height=150,
            width=500,
            bg="#FCFFEE"
        )
        self.home_foot_frame.grid(column=0, row=4)

        # ファイル読み込みの作成と配置
        self.input_button = tkinter.Button(
            self.home_input_frame,
            image = self.input_button_image,
            highlightbackground='#FCFFEE'
        )
        self.input_button.pack()

        # 使い方の作成と配置
        self.using_button = tkinter.Button(
            self.home_operation_frame,
            image = self.using_button_image,
            highlightbackground='#FCFFEE'
        )
        self.using_button.pack(fill = 'x', padx=25, side = 'left')

        # 実行の作成と配置
        self.done_button = tkinter.Button(
            self.home_operation_frame,
            image = self.done_button_image,
            highlightbackground='#FCFFEE'
        )
        self.done_button.pack(fill = 'x', padx=25, side = 'right')
        return

    def create_end_frame(self):
        # 画像の読み込み
        complete_image_ = Image.open('./app/image/完了_M.png')
        self.complete_image = ImageTk.PhotoImage(complete_image_)

        # キャンバスの配置
        self.complete_canvas = tkinter.Canvas(
            self.end_frame,
            height=500,
            width=700,
            highlightbackground='#FCFFEE',
            bg="#FCFFEE"
        )
        self.complete_canvas.create_image(
            30, 30,
            image=self.complete_image,
            anchor=tkinter.NW,
            tag="image",
        )
        self.complete_canvas.grid(column=0, row=0)
        # 画像の配置


        return

    def select_open_file(self, file_types):
        'オープンするファイル選択画面を表示'

        # ファイル選択ダイアログを表示
        file_path = tkinter.filedialog.askopenfilename(
            initialdir=".",
            filetypes=file_types
        )
        return file_path

    # def make_draw_image(self, name, index):
    #     return lambda: self.draw_image(self, "", index)

    def draw_image(self, name, index):
        i = index
        alli = self.model.nows[index].get()
        'buttonの切り替え'
        state_ = self.model.frame_state[alli].get()
        if state_ == 1: # 選択されているので、削除ボタンを表示
            self.state_button2[i].tkraise()
        else: # 選択されていないので、追加ボタンを表示
            self.state_button[i].tkraise()
        
        '画像をキャンバスに描画'
        self.frame[i].delete('all')
        if alli < 0 or len(self.model.frames) <= alli:
            return
        image = self.model.frames[alli]
        self.frame[i].create_image(
            0, 0,
            image=image,
            anchor=tkinter.NW,
            tag="image"
        )



class Controller():

    def __init__(self, app, model, view):
        self.master = app
        self.model = model
        self.view = view

        self.set_events()
        self.model.mode.set(0)
        self.model.now.set(0)


    def set_events(self):
        '受け付けるイベントを設定する'

        # 動画選択ボタン押し下げイベント受付
        self.view.input_button['command'] = self.push_load_button

        # モノクロON/OFFボタン押し下げイベント受付
        self.view.next_frame_button['command'] = self.push_next_frame_button

        # フリップON/OFFボタン押し下げイベント受付
        self.view.prev_frame_button['command'] = self.push_prev_frame_button

        # 編集とプレビューの切り替えイベント受付
        self.view.mode_button['command'] = self.push_mode_button
        self.view.mode_button2['command'] = self.push_mode_button

        for x in range(5):
            self.view.state_button[x]['command'] = self.make_push_state_button(x)    
        
        for x in range(5):
            self.view.state_button2[x]['command'] = self.make_push_state_button(x)    

        # ホームからプレビューへの画面遷移
        self.view.done_button['command'] = self.push_done_button

        # プレビューからホームへの画面遷移
        self.view.to_home_button['command'] = self.push_to_home_button

        # プレビューからエンドへの画面遷移
        self.view.making_pptx_button['command'] = self.push_making_pptx_button

    def push_load_button(self):
        '動画選択ボタンが押された時の処理'

        file_types = [
            ("MOVファイル", "*.mov"),
            ("MP4ファイル", "*.mp4"),
        ]

        # ファイル選択画面表示
        file_path = self.view.select_open_file(file_types)
        if len(file_path) != 0:
            # 動画オブジェクト生成
            self.model.create_video(file_path)

    def push_next_frame_button(self):
        self.model.next_frame()

    def push_prev_frame_button(self):
        self.model.previous_frame()

    # ここではボタン押したら、今の数字のやつのstate切り替えと、再び表示がしたい、真ん中に真ん中を入れたらおけ
    def make_push_state_button(self, x):
        return lambda: self.push_state_button(x)

    def push_state_button(self, x):
        index = self.model.nows[x].get()
        if index < 0 or len(self.model.frames) <= index:
            return
        self.model.frame_state[index].set((self.model.frame_state[index].get()+1)%2)
        # print([self.model.frame_state[x].get() for x in range(len(self.model.frame_state))])
        # 全体の更新
        self.model.now.set(self.model.now.get())

    def push_mode_button(self):
        self.model.mode.set((self.model.mode.get()+1)%2)
        if self.model.mode.get() == 0: # プレビュー
            self.view.head_canvas.delete('all')
            self.view.head_canvas.create_image(
                30, 30,
                image=self.view.preview_top_image,
                anchor=tkinter.NW,
                tag="image",
            )
            self.view.making_pptx_button.tkraise()
            self.view.to_home_button.tkraise()
            self.view.mode_button.tkraise()
            self.model.now.set(0)
        else: # 編集モード
            self.view.head_canvas.delete('all')
            self.view.head_canvas.create_image(
                30, 31,
                image=self.view.edit_top_image,
                anchor=tkinter.NW,
                tag="image",
            )
            self.view.mode_button2.tkraise()
            self.view.to_home_button.lower()
            self.view.mode_button.lower()

        
        # ここに編集モードとプレビューモードのviewの変更をかく
        self.model.set_nows("", "", "")

    def push_done_button(self):
        self.view.main_frame.tkraise()

    def push_to_home_button(self):
        self.view.home_frame.tkraise()

    def push_making_pptx_button(self):
        self.view.end_frame.tkraise()
        self.making_pptx()

    def making_pptx(self):
        with tempfile.TemporaryDirectory() as dname:
            # tempdirに画像を保存
            for (index, img) in tqdm(enumerate(self.model.perfect_frames)):
                if self.model.frame_state[index].get()==0:
                    continue
                cv2.imwrite(f'{dname}/{index:04}.jpeg', img)
            # 保存した名前を取得
            img_names = os.listdir(dname)
            img_names = [name for name in img_names if name.endswith(".jpeg")]
            img_names.sort()
            
            # imanishi_pptx
            self.imanishi_pptx(dname, img_names)

    def put_pic(self, slide, path, pic_left, pic_top, pic_width, pic_height):
        #画像追加
        image = slide.shapes.add_picture(path, pic_left, pic_top, pic_width, pic_height) 
        image.line.color.rgb = RGBColor(0, 0, 0)
        image.line.width = Pt(1.5)

    def put_text(self, slide, pic_left, pic_top, str, size):
        #テキストを追加
        textbox = slide.shapes.add_textbox(pic_left, pic_top, Pt(size), Pt(size))
        tf = textbox.text_frame
        tf.text = str
        tf.paragraphs[0].font.size = Pt(size)  # font size
        tf.paragraphs[0].alignment = PP_ALIGN.RIGHT

    def put_arrow(self, slide, pic_left2, pic_top2, size):
        #矢印出力
        rect0 = slide.shapes.add_shape(		# shapeオブジェクト➀を追加
                MSO_SHAPE.RIGHT_ARROW,   	                    # 図形の種類を[丸角四角形]に指定
                pic_left2, pic_top2,               # 挿入位置の指定：左からの座標と上からの座標の指定
                size, size)               # 挿入図形の幅と高さの指定
        rect0.fill.solid()                                   # shapeオブジェクト➀を単色で塗り潰す
        rect0.fill.fore_color.rgb = RGBColor(74, 126, 187)  # RGB指定で色を指定

    def put_sign(self,slide,  SIGN_DIR, sign_names):
        pic_top3 = Cm(0)
        width = Cm(5.49)
        height = Cm(3.94)
        for i,name in enumerate(sign_names):
            if(i < 2):
                slide.shapes.add_picture(SIGN_DIR+name, Cm(-6.5), pic_top3, height=height) 
                pic_top3 += height
            elif i == 2:
                pic_top3 = Cm(-1.05)
                width = Cm(5.92)
                height = Cm(4.92)
                slide.shapes.add_picture(SIGN_DIR+name, Cm(29.69), pic_top3, height=height) 
                #pic_top3 += height
            elif i == 3:
                slide.shapes.add_picture(SIGN_DIR+name, Cm(31.69), Cm(2.8), height=height) 
                #pic_top3 += height
            elif i == 4:
                height = Cm(4.92)
                slide.shapes.add_picture(SIGN_DIR+name, Cm(30.96),Cm(5.9), height=height) 
            else:
                #width = Cm(4.92)
                height = Cm(5.61)
                slide.shapes.add_picture(SIGN_DIR+name, Cm(30.96), Cm(11.12), height=height) 

    def imanishi_pptx(self, IMG_DIR, img_names):
        IMG_DIR = IMG_DIR + "/"
        #画像の格納ディレクトリ
        SIGN_DIR = "./app/sign/"
        #img画像のファイル名を取得
        sign_names = os.listdir(SIGN_DIR)
        sign_names = [name for name in sign_names if name.endswith(".png")]
        sign_names.sort()#昇順にsort

        prs = pptx.Presentation()
        prs.slide_width = Inches(11.69) #A4サイズ
        prs.slide_height = Inches(8.27)
        slide_width = prs.slide_width
        slide_height = prs.slide_height

        #画像のアスペクト比を取得
        im = Image.open(IMG_DIR + '/' +img_names[0])
        aspect_ratio = im.width /im.height

        #画像の高さの設定と幅の取得
        pic_height = Cm(9.5)
        pic_width = aspect_ratio * pic_height

        # タイトルスライド
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        textbox = slide.shapes.add_textbox(slide_width/2-Inches(5)/2, slide_height/2-Inches(1)/2, Inches(5), Inches(1))
        tf = textbox.text_frame
        tf.text = "タイトルを入力"
        tf.paragraphs[0].font.size = Pt(50)  # font size
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

                #画像を１枚のパワポに出力 1段4枚ずつ
        pic_left = Cm(2.5)
        pic_left2 = Cm(7.2)
        for i, name in enumerate(img_names):
            path = IMG_DIR + name

            if i % 8 == 0:
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                pic_top = Cm(0.5)

            #画像追加
            self.put_pic( slide, path, pic_left, pic_top, pic_width, pic_height)
            #テキストを追加
            self.put_text(slide,pic_left-Cm(1), pic_top, str(i+1), 28)

            pic_left += Cm(7)

            if i % 4 == 3:
                pic_top += Cm(10)
                pic_left = Cm(2.5)
                pic_left2 = Cm(7.2)
                pic_top2 += Cm(11)
            elif i != len(img_names)-1:
            #矢印出力
                pic_top2 = pic_top + pic_height/2 - Cm(1)
                self.put_arrow( slide, pic_left2, pic_top2, Cm(2))
                pic_left2 += Cm(7) 

        #画像を２枚ずつパワポに出力
        pic_height = Cm(16)
        pic_width = aspect_ratio * pic_height
        pic_top = ( slide_height - pic_height ) / 2

        # 連番で2枚ずつのスライドを作る疑似コード
        pre_path = None
        for i,name in enumerate(img_names):
            path = IMG_DIR + name
            if i==0:
                pre_path = path
                continue
            #スライドを増やす
            slide = prs.slides.add_slide(prs.slide_layouts[6]) 
            #pre_pathの画像を←に配置
            pic_left = ( slide_width/2 - pic_width ) / 2
            self.put_pic( slide, pre_path, pic_left, pic_top, pic_width, pic_height)
            self.put_text( slide, pic_left-Cm(1.5), pic_top, str(i), 36)
            pre_path = path
            #pathの画像を→に配置
            pic_left += slide_width/2
            self.put_pic( slide, path, pic_left, pic_top, pic_width, pic_height)
            #矢印を追加
            ratio = 0.45
            pic_left2 = slide_width/2 - pic_width*ratio/2
            pic_top2 = slide_height/2 - pic_width*ratio/2
            self.put_arrow( slide, pic_left2, pic_top2, pic_width*ratio)
            #手の写真を追加
            self.put_sign( slide, SIGN_DIR, sign_names)
            #テキストを追加
            self.put_text( slide, pic_left-Cm(1.5), pic_top, str(i+1), 36)
            pre_path = path

        prs.save("./らくらくトリセツ.pptx") 

app = tkinter.Tk()

app.title("らくらくトリセツ")

model = Model()
view = View(app, model)
controller = Controller(app, model, view)

app.mainloop()